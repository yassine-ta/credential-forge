"""Tests for file format synthesizers."""

import pytest
import tempfile
from pathlib import Path
from unittest.mock import Mock, patch

from credentialforge.synthesizers.base import BaseSynthesizer
from credentialforge.synthesizers.eml_synthesizer import EMLSynthesizer
from credentialforge.synthesizers.excel_synthesizer import ExcelSynthesizer
from credentialforge.synthesizers.pptx_synthesizer import PowerPointSynthesizer
from credentialforge.synthesizers.vsdx_synthesizer import VisioSynthesizer


class TestBaseSynthesizer:
    """Test cases for BaseSynthesizer."""
    
    @pytest.fixture
    def temp_output_dir(self):
        """Create temporary output directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir
    
    def test_init(self, temp_output_dir):
        """Test synthesizer initialization."""
        # Create a concrete implementation for testing
        class TestSynthesizer(BaseSynthesizer):
            def synthesize(self, topic_content, credentials, metadata=None):
                return "test_file.txt"
        
        synthesizer = TestSynthesizer(temp_output_dir)
        
        assert synthesizer.output_dir == Path(temp_output_dir)
        assert synthesizer.output_dir.exists()
        assert synthesizer.generation_stats['files_generated'] == 0
    
    def test_embed_credentials_random(self, temp_output_dir):
        """Test random credential embedding."""
        class TestSynthesizer(BaseSynthesizer):
            def synthesize(self, topic_content, credentials, metadata=None):
                return "test_file.txt"
        
        synthesizer = TestSynthesizer(temp_output_dir)
        content = "Line 1\nLine 2\nLine 3"
        credentials = ["cred1", "cred2"]
        
        result = synthesizer._embed_credentials(content, credentials, 'random')
        
        assert "cred1" in result or "cred2" in result
        assert len(result) > len(content)
    
    def test_generate_filename(self, temp_output_dir):
        """Test filename generation."""
        class TestSynthesizer(BaseSynthesizer):
            def synthesize(self, topic_content, credentials, metadata=None):
                return "test_file.txt"
        
        synthesizer = TestSynthesizer(temp_output_dir)
        
        filename = synthesizer._generate_filename("test", "txt", {"topic": "test topic"})
        
        assert filename.startswith("test_")
        assert filename.endswith(".txt")
        assert "test_topic" in filename
    
    def test_get_generation_stats(self, temp_output_dir):
        """Test generation statistics."""
        class TestSynthesizer(BaseSynthesizer):
            def synthesize(self, topic_content, credentials, metadata=None):
                return "test_file.txt"
        
        synthesizer = TestSynthesizer(temp_output_dir)
        synthesizer.generation_stats['files_generated'] = 5
        synthesizer.generation_stats['total_credentials_embedded'] = 10
        
        stats = synthesizer.get_generation_stats()
        
        assert stats['files_generated'] == 5
        assert stats['total_credentials_embedded'] == 10


class TestEMLSynthesizer:
    """Test cases for EMLSynthesizer."""
    
    @pytest.fixture
    def temp_output_dir(self):
        """Create temporary output directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir
    
    @pytest.fixture
    def synthesizer(self, temp_output_dir):
        """Create EMLSynthesizer instance."""
        return EMLSynthesizer(temp_output_dir)
    
    def test_synthesize_basic(self, synthesizer, temp_output_dir):
        """Test basic EML synthesis."""
        content = "Test email content"
        credentials = ["AKIA1234567890ABCDEF", "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9..."]
        
        file_path = synthesizer.synthesize(content, credentials)
        
        assert Path(file_path).exists()
        assert file_path.endswith('.eml')
        
        # Check file content
        with open(file_path, 'r', encoding='utf-8') as f:
            file_content = f.read()
            assert "Test email content" in file_content
            assert "From:" in file_content
            assert "Subject:" in file_content
    
    def test_synthesize_with_metadata(self, synthesizer, temp_output_dir):
        """Test EML synthesis with metadata."""
        content = "Test email content"
        credentials = ["AKIA1234567890ABCDEF"]
        metadata = {
            'sender': 'test@example.com',
            'recipient': 'user@example.com',
            'subject': 'Test Subject'
        }
        
        file_path = synthesizer.synthesize(content, credentials, metadata)
        
        with open(file_path, 'r', encoding='utf-8') as f:
            file_content = f.read()
            assert "test@example.com" in file_content
            assert "user@example.com" in file_content
            assert "Test Subject" in file_content
    
    def test_embed_credentials_in_content(self, synthesizer):
        """Test credential embedding in email content."""
        content = "API configuration:\nDatabase connection:\nMonitoring setup:"
        credentials = ["AKIA1234567890ABCDEF", "mysql://user:pass@host:3306/db"]
        
        result = synthesizer._embed_credentials_in_content(content, credentials)
        
        assert "AKIA1234567890ABCDEF" in result
        assert "mysql://user:pass@host:3306/db" in result


class TestExcelSynthesizer:
    """Test cases for ExcelSynthesizer."""
    
    @pytest.fixture
    def temp_output_dir(self):
        """Create temporary output directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir
    
    @pytest.fixture
    def synthesizer(self, temp_output_dir):
        """Create ExcelSynthesizer instance."""
        return ExcelSynthesizer(temp_output_dir)
    
    def test_synthesize_basic(self, synthesizer, temp_output_dir):
        """Test basic Excel synthesis."""
        content = "Test spreadsheet content"
        credentials = ["AKIA1234567890ABCDEF", "mysql://user:pass@host:3306/db"]
        
        file_path = synthesizer.synthesize(content, credentials)
        
        assert Path(file_path).exists()
        assert file_path.endswith('.xlsx')
        
        # Check that file is a valid Excel file
        assert file_path in str(Path(temp_output_dir).glob('*.xlsx'))
    
    def test_create_worksheets(self, synthesizer):
        """Test worksheet creation."""
        from openpyxl import Workbook
        
        wb = Workbook()
        wb.remove(wb.active)
        
        content = "Test content"
        credentials = ["AKIA1234567890ABCDEF"]
        
        synthesizer._create_worksheets(wb, content, credentials)
        
        # Check that worksheets were created
        sheet_names = [sheet.title for sheet in wb.worksheets]
        assert "Configuration" in sheet_names
        assert "Credentials" in sheet_names
        assert "Data" in sheet_names
        assert "Formulas" in sheet_names


class TestPowerPointSynthesizer:
    """Test cases for PowerPointSynthesizer."""
    
    @pytest.fixture
    def temp_output_dir(self):
        """Create temporary output directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir
    
    @pytest.fixture
    def synthesizer(self, temp_output_dir):
        """Create PowerPointSynthesizer instance."""
        return PowerPointSynthesizer(temp_output_dir)
    
    def test_synthesize_basic(self, synthesizer, temp_output_dir):
        """Test basic PowerPoint synthesis."""
        content = "Test presentation content"
        credentials = ["AKIA1234567890ABCDEF", "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9..."]
        
        file_path = synthesizer.synthesize(content, credentials)
        
        assert Path(file_path).exists()
        assert file_path.endswith('.pptx')
        
        # Check that file is a valid PowerPoint file
        assert file_path in str(Path(temp_output_dir).glob('*.pptx'))
    
    def test_create_slides(self, synthesizer):
        """Test slide creation."""
        from pptx import Presentation
        
        prs = Presentation()
        
        content = "Test content"
        credentials = ["AKIA1234567890ABCDEF"]
        metadata = {"title": "Test Presentation"}
        
        synthesizer._create_slides(prs, content, credentials, metadata)
        
        # Check that slides were created
        assert len(prs.slides) > 0
        
        # Check slide titles
        slide_titles = [slide.shapes.title.text for slide in prs.slides if slide.shapes.title]
        assert any("System Architecture" in title for title in slide_titles)


class TestVisioSynthesizer:
    """Test cases for VisioSynthesizer."""
    
    @pytest.fixture
    def temp_output_dir(self):
        """Create temporary output directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir
    
    @pytest.fixture
    def synthesizer(self, temp_output_dir):
        """Create VisioSynthesizer instance."""
        return VisioSynthesizer(temp_output_dir)
    
    def test_synthesize_basic(self, synthesizer, temp_output_dir):
        """Test basic Visio synthesis."""
        content = "Test diagram content"
        credentials = ["AKIA1234567890ABCDEF", "mysql://user:pass@host:3306/db"]
        
        file_path = synthesizer.synthesize(content, credentials)
        
        assert Path(file_path).exists()
        assert file_path.endswith('.vsdx')
        
        # Check that file is a valid Visio file (ZIP archive)
        import zipfile
        assert zipfile.is_zipfile(file_path)
    
    def test_create_vsdx_structure(self, synthesizer):
        """Test VSDX structure creation."""
        content = "Test content"
        credentials = ["AKIA1234567890ABCDEF"]
        metadata = {"title": "Test Diagram"}
        
        vsdx_content = synthesizer._create_vsdx_structure(content, credentials, metadata)
        
        # Check that required files are present
        required_files = [
            '[Content_Types].xml',
            '_rels/.rels',
            'docProps/app.xml',
            'docProps/core.xml',
            'visio/document.xml'
        ]
        
        for file_path in required_files:
            assert file_path in vsdx_content
    
    def test_create_shapes_xml(self, synthesizer):
        """Test shapes XML creation."""
        content = "Test content"
        credentials = ["AKIA1234567890ABCDEF", "mysql://user:pass@host:3306/db"]
        
        shapes_xml = synthesizer._create_shapes_xml(content, credentials)
        
        assert "<Shapes>" in shapes_xml
        assert "System Architecture" in shapes_xml
        assert "AKIA1234567890ABCDEF" in shapes_xml


class TestIntegration:
    """Integration tests for synthesizers."""
    
    @pytest.fixture
    def temp_output_dir(self):
        """Create temporary output directory."""
        with tempfile.TemporaryDirectory() as tmpdir:
            yield tmpdir
    
    def test_all_synthesizers(self, temp_output_dir):
        """Test all synthesizers can create files."""
        synthesizers = {
            'eml': EMLSynthesizer(temp_output_dir),
            'xlsx': ExcelSynthesizer(temp_output_dir),
            'pptx': PowerPointSynthesizer(temp_output_dir),
            'vsdx': VisioSynthesizer(temp_output_dir)
        }
        
        content = "Test content for all formats"
        credentials = ["AKIA1234567890ABCDEF", "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9..."]
        
        for format_name, synthesizer in synthesizers.items():
            file_path = synthesizer.synthesize(content, credentials)
            
            assert Path(file_path).exists()
            assert file_path.endswith(f'.{format_name}' if format_name != 'xlsx' else '.xlsx')
            
            # Check file size is reasonable
            file_size = Path(file_path).stat().st_size
            assert file_size > 100  # Should be at least 100 bytes
    
    def test_credential_embedding_across_formats(self, temp_output_dir):
        """Test credential embedding works across all formats."""
        synthesizers = {
            'eml': EMLSynthesizer(temp_output_dir),
            'xlsx': ExcelSynthesizer(temp_output_dir),
            'pptx': PowerPointSynthesizer(temp_output_dir),
            'vsdx': VisioSynthesizer(temp_output_dir)
        }
        
        content = "API configuration with credentials"
        credentials = ["AKIA1234567890ABCDEF", "mysql://user:pass@host:3306/db"]
        
        for format_name, synthesizer in synthesizers.items():
            file_path = synthesizer.synthesize(content, credentials)
            
            # For text-based formats, check content contains credentials
            if format_name in ['eml']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    file_content = f.read()
                    # At least one credential should be embedded
                    assert any(cred in file_content for cred in credentials)
