"""PPTX format synthesizer using agent-generated content."""

import random
from pathlib import Path
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .format_synthesizer import FormatSynthesizer
from ..utils.exceptions import SynthesizerError


class PPTXFormatSynthesizer(FormatSynthesizer):
    """PPTX format synthesizer that structures agent-generated content."""
    
    def __init__(self, output_dir: str, ultra_fast_mode: bool = False):
        """Initialize PPTX format synthesizer.
        
        Args:
            output_dir: Output directory for generated files
            ultra_fast_mode: Enable ultra-fast mode with minimal validation
        """
        super().__init__(output_dir, ultra_fast_mode)
        self.format_type = 'pptx'
    
    def synthesize(self, content_structure: Dict[str, Any]) -> str:
        """Structure content into PPTX format.
        
        Args:
            content_structure: Generated content from agents
            
        Returns:
            Path to generated PPTX file
        """
        try:
            # Validate content structure
            self._validate_content_structure(content_structure)
            
            # Create presentation
            prs = Presentation()
            
            # Create slides from sections
            self._create_slides_from_sections(prs, content_structure)
            
            # Generate filename and save
            filename = self._generate_filename(content_structure)
            file_path = self._get_file_path(filename)
            
            # Save presentation
            prs.save(str(file_path))
            
            # Log stats
            self._log_generation_stats(content_structure)
            
            return str(file_path)
            
        except Exception as e:
            self.generation_stats['errors'] += 1
            raise SynthesizerError(f"PPTX synthesis failed: {e}")
    
    def _create_slides_from_sections(self, prs: Presentation, content_structure: Dict[str, Any]) -> None:
        """Create slides from content sections."""
        sections = content_structure.get('sections', [])
        credentials = content_structure.get('credentials', [])
        language = content_structure.get('language', 'en')
        
        # Title slide
        self._create_title_slide(prs, content_structure)
        
        # Content slides
        for i, section in enumerate(sections):
            if i == 0:  # Skip title section as it's handled separately
                continue
            
            self._create_content_slide(prs, section, language)
        
        # Credentials slide (if credentials exist)
        if credentials:
            self._create_credentials_slide(prs, credentials, language)
    
    def _create_title_slide(self, prs: Presentation, content_structure: Dict[str, Any]) -> None:
        """Create title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = content_structure.get('title', 'Document Title')
        
        # Set subtitle if available
        subtitle = slide.placeholders[1]
        metadata = content_structure.get('metadata', {})
        topic = metadata.get('topic', 'System Documentation')
        subtitle.text = f"Topic: {topic}"
        
        # Style the title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
    
    def _create_content_slide(self, prs: Presentation, section: Dict[str, str], language: str) -> None:
        """Create content slide from section."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = section.get('title', 'Section')
        
        # Set content
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()
        
        content = section.get('content', '')
        
        # Split content into paragraphs
        paragraphs = content.split('\n\n')
        
        for i, paragraph in enumerate(paragraphs):
            if paragraph.strip():
                p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
                p.text = paragraph.strip()
                p.font.size = Pt(18)
                
                # Make first paragraph bold
                if i == 0:
                    p.font.bold = True
    
    def _create_credentials_slide(self, prs: Presentation, credentials: list, language: str) -> None:
        """Create credentials slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = self._get_credentials_slide_title(language)
        
        # Set content
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()
        
        # Add credentials
        for i, cred in enumerate(credentials):
            label = cred.get('label', cred.get('type', 'Credential'))
            value = cred.get('value', '')
            
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = f"{label}: {value}"
            p.font.size = Pt(20)
            p.font.bold = True
            
            # Color code credentials
            if 'password' in cred.get('type', '').lower():
                p.font.color.rgb = RGBColor(255, 0, 0)  # Red for passwords
            elif 'api' in cred.get('type', '').lower():
                p.font.color.rgb = RGBColor(0, 0, 255)  # Blue for API keys
            else:
                p.font.color.rgb = RGBColor(0, 128, 0)  # Green for others
        
        # Add note in speaker notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text = f"{self._get_credentials_slide_title(language)}:\n\n"
        
        for cred in credentials:
            label = cred.get('label', cred.get('type', 'Credential'))
            value = cred.get('value', '')
            notes_text += f"{label}: {value}\n"
        
        notes_text_frame.text = notes_text
    
    def _get_credentials_slide_title(self, language: str) -> str:
        """Get localized credentials slide title."""
        titles = {
            'en': 'Configuration Credentials',
            'fr': 'Identifiants de Configuration',
            'es': 'Credenciales de Configuración',
            'de': 'Konfigurationsanmeldedaten',
            'it': 'Credenziali di Configurazione'
        }
        
        return titles.get(language, titles['en'])
    
    def _generate_filename(self, content_structure: Dict[str, Any]) -> str:
        """Generate PPTX filename."""
        title = content_structure.get('title', 'presentation')
        timestamp = self._get_current_timestamp()
        random_id = random.randint(1000, 9999)
        
        # Clean title for filename
        clean_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        clean_title = clean_title.replace(' ', '_').lower()
        
        return f"presentation_{clean_title}_{timestamp}_{random_id}.pptx"
    
    def _get_current_timestamp(self) -> str:
        """Get current timestamp for filename."""
        from datetime import datetime
        return datetime.now().strftime('%Y%m%d_%H%M%S')
